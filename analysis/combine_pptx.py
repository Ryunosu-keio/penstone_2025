import collections.abc
import collections.abc # インポートしないとエラーが発生する
from pptx import Presentation # プレゼンテーションを作成
from pptx.util import Inches  # インチ
from pptx.enum.text import PP_ALIGN  # 中央揃えにする用
from pptx.util import Cm, Pt # センチ、ポイント
from pptx import Presentation

def copy_slide_from_to(src_presentation, tgt_presentation):
    for slide in src_presentation.slides:
        slide_id = tgt_presentation.slides.add_slide(slide.slide_layout).slide_id
        src_slide = src_presentation.slides.get(slide_id)
        tgt_presentation.slides._sldIdLst.remove(tgt_presentation.slides._sldIdLst[-1])
        tgt_presentation.slides._sldIdLst.insert(-1, src_slide.element)

# 2つのPowerPointファイルを読み込みます
pptx_file1 = 'path_to_first_presentation.pptx'
pptx_file2 = 'path_to_second_presentation.pptx'

# 新しいプレゼンテーションを作成します
combined_presentation = Presentation()

# 最初のプレゼンテーションからスライドをコピーします
src_presentation1 = Presentation(pptx_file1)
copy_slide_from_to(src_presentation1, combined_presentation)

# 2番目のプレゼンテーションからスライドをコピーします
src_presentation2 = Presentation(pptx_file2)
copy_slide_from_to(src_presentation2, combined_presentation)

# 新しいファイルとして保存します
combined_presentation.save('combined_presentation.pptx')
